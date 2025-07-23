import json
from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.enum.shapes import MSO_SHAPE_TYPE


def pptx_to_json(pptx_path) -> list[dict]:
    def extract_text_from_shape(shape, slide_num, element_counter):
        elements = []
        if shape.has_text_frame:
            full_text = shape.text.strip()
            if full_text:
                font_name = "default"
                font_size = 12.0
                bold = False
                italic = False

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            font = run.font
                            font_name = font.name if font.name else "default"
                            font_size = font.size.pt if font.size else 12.0
                            bold = font.bold if font.bold is not None else False
                            italic = font.italic if font.italic is not None else False
                            break
                    if font_name != "default":
                        break

                element_info = {
                    "id_element": f"s{slide_num}_e{element_counter[0]}",
                    "type": "text",
                    "text": full_text,
                    "font_name": font_name,
                    "font_size": font_size,
                    "bold": bold,
                    "italic": italic,
                    "length": len(full_text)
                }
                elements.append(element_info)
                element_counter[0] += 1

        return elements

    def extract_image_from_shape(shape, slide_num, element_counter):
        elements = []
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                width = shape.width.pt if shape.width else 0
                height = shape.height.pt if shape.height else 0

                element_info = {
                    "id_element": f"s{slide_num}_e{element_counter[0]}",
                    "type": "image",
                    "width": round(width, 2),
                    "height": round(height, 2)
                }
                elements.append(element_info)
                element_counter[0] += 1

        except Exception:
            pass

        return elements

    def extract_text_from_table(table_shape, slide_num, element_counter):
        elements = []
        try:
            table = table_shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        element_info = {
                            "id_element": f"s{slide_num}_e{element_counter[0]}",
                            "type": "table_cell",
                            "text": cell.text.strip(),
                            "font_name": "default",
                            "font_size": 12.0,
                            "bold": False,
                            "italic": False,
                            "length": len(cell.text.strip())
                        }
                        elements.append(element_info)
                        element_counter[0] += 1
        except Exception:
            pass

        return elements

    def extract_from_group(group, slide_num, element_counter):
        elements = []
        for shape in group.shapes:
            if isinstance(shape, GroupShape):
                elements.extend(extract_from_group(shape, slide_num, element_counter))
            elif hasattr(shape, 'table') and shape.table:
                elements.extend(extract_text_from_table(shape, slide_num, element_counter))
            elif shape.has_text_frame:
                elements.extend(extract_text_from_shape(shape, slide_num, element_counter))
            else:
                elements.extend(extract_image_from_shape(shape, slide_num, element_counter))
        return elements

    try:
        prs = Presentation(pptx_path)
        data = []

        for i, slide in enumerate(prs.slides, start=1):
            slide_info = {
                "slide_number": i,
                "elements": []
            }

            element_counter = [1]

            for shape in slide.shapes:
                try:
                    if isinstance(shape, GroupShape):
                        elements = extract_from_group(shape, i, element_counter)
                        slide_info["elements"].extend(elements)

                    elif hasattr(shape, 'table') and shape.table:
                        elements = extract_text_from_table(shape, i, element_counter)
                        slide_info["elements"].extend(elements)

                    elif shape.has_text_frame:
                        elements = extract_text_from_shape(shape, i, element_counter)
                        slide_info["elements"].extend(elements)

                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        elements = extract_image_from_shape(shape, i, element_counter)
                        slide_info["elements"].extend(elements)

                    elif hasattr(shape, 'text') and shape.text.strip():
                        element_info = {
                            "id_element": f"s{i}_e{element_counter[0]}",
                            "type": "text",
                            "text": shape.text.strip(),
                            "font_name": "default",
                            "font_size": 12.0,
                            "bold": False,
                            "italic": False,
                            "length": len(shape.text.strip())
                        }
                        slide_info["elements"].append(element_info)
                        element_counter[0] += 1

                except Exception:
                    continue

            data.append(slide_info)

        return data

    except Exception as e:
        return [{"error": f"Ошибка при обработке файла: {str(e)}"}]

import json
from pptx import Presentation
from pptx.shapes.group import GroupShape


def load_json_structure(json_path):
    with open(json_path, 'r', encoding='utf-8') as f:
        return json.load(f)

def extract_text_elements_from_shape(shape, elements_list):
    if shape.has_text_frame and shape.text.strip():
        elements_list.append(shape)

def extract_text_elements_from_table(table_shape, elements_list):
    try:
        table = table_shape.table
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    elements_list.append(cell)
    except Exception as e:
        print(f"Ошибка: {e}")

def extract_text_elements_from_group(group, elements_list):
    for shape in group.shapes:
        if isinstance(shape, GroupShape):
            extract_text_elements_from_group(shape, elements_list)
        elif hasattr(shape, 'table') and shape.table:
            extract_text_elements_from_table(shape, elements_list)
        elif shape.has_text_frame and shape.text.strip():
            elements_list.append(shape)

def extract_slide_text_elements(slide):
    elements_list = []

    for shape in slide.shapes:
        try:
            if isinstance(shape, GroupShape):
                extract_text_elements_from_group(shape, elements_list)
            elif hasattr(shape, 'table') and shape.table:
                extract_text_elements_from_table(shape, elements_list)
            elif shape.has_text_frame and shape.text.strip():
                elements_list.append(shape)
            elif hasattr(shape, 'text') and shape.text.strip():
                elements_list.append(shape)
        except Exception as e:
            continue

    return elements_list

def replace_shape_text(shape_or_cell, new_text):
    try:
        if hasattr(shape_or_cell, 'text_frame'):
            for paragraph in shape_or_cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        run.text = new_text
                        return
        elif hasattr(shape_or_cell, 'text'):
            shape_or_cell.text = new_text
    except Exception as e:
        print(f"Ошибка: {e}")


def replace_text_with_json_structure(input_pptx_path, json_structure):
    try:
        slides_data = json_structure
        prs = Presentation(input_pptx_path)

        for slide_index, slide in enumerate(prs.slides):
            if slide_index >= len(slides_data):
                break

            slide_data = slides_data[slide_index]
            json_elements = slide_data['elements']

            slide_text_elements = extract_slide_text_elements(slide)

            replacements_made = 0
            for i, (slide_element, json_element) in enumerate(zip(slide_text_elements, json_elements)):
                try:
                    old_text = getattr(slide_element, 'text', 'N/A')
                    new_text = json_element['text']

                    replace_shape_text(slide_element, new_text)
                    replacements_made += 1

                except Exception as e:
                    continue
        return prs

    except Exception as e:
        print(f"Ошибка: {e}")